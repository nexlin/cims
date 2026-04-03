#!/usr/bin/env python3
"""CIMS 시스템 기술 문서 PPT 생성 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(37, 99, 235)
DARK = RGBColor(26, 29, 46)
GRAY = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(219, 234, 254)
LIGHT_GREEN = RGBColor(220, 252, 231)
LIGHT_GRAY = RGBColor(243, 244, 246)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)

# ── 헬퍼 ─────────────────────────────────────────────────────

def _title_bar(slide, title):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    r.fill.solid(); r.fill.fore_color.rgb = BLUE; r.line.fill.background()
    tf = r.text_frame; tf.margin_top = Pt(10); tf.margin_left = Pt(24)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

def title_slide(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BLUE
    tx = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(2.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(200,220,255); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(20)

def section_slide(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
    tx = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(16); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)

def _box(slide, x, y, w, h, text, sub, fill_rgb, text_rgb=DARK):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill_rgb; r.line.color.rgb = BLUE; r.line.width = Pt(1.5)
    r.shadow.inherit = False
    tf = r.text_frame; tf.word_wrap = True; tf.margin_top = Pt(6); tf.margin_left = Pt(8)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = text_rgb; p.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(10); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER

def _arrow_text(slide, x, y, w, text):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    p = tx.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(9); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER

def content_slide(title, items):
    """items: list of (heading, details_list) tuples"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(s, title)
    y = 1.1
    for heading, details in items:
        tx = s.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12), Inches(0.35))
        p = tx.text_frame.paragraphs[0]; p.text = heading; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BLUE
        y += 0.38
        for d in details:
            tx2 = s.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11.5), Inches(0.28))
            p2 = tx2.text_frame.paragraphs[0]; p2.text = f"• {d}"; p2.font.size = Pt(14); p2.font.color.rgb = DARK
            y += 0.3
        y += 0.1

def code_slide(title, code_text, note=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(s, title)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.5))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(30,30,46); r.line.fill.background()
    tf = r.text_frame; tf.word_wrap = True; tf.margin_top = Pt(12); tf.margin_left = Pt(16)
    for i, line in enumerate(code_text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12); p.font.name = 'Consolas'; p.font.color.rgb = RGBColor(205,214,244)
        p.space_before = Pt(1); p.space_after = Pt(1)
    if note:
        tx = s.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
        p = tx.text_frame.paragraphs[0]; p.text = note; p.font.size = Pt(11); p.font.color.rgb = GRAY; p.font.italic = True

def table_slide(title, headers, rows, note=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(s, title)
    cols = len(headers); r_count = len(rows) + 1
    row_h = min(0.42, 5.5 / r_count)
    tbl = s.shapes.add_table(r_count, cols, Inches(0.5), Inches(1.1), Inches(12.3), Inches(row_h * r_count)).table
    for i, h in enumerate(headers):
        c = tbl.cell(0, i); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = DARK
        for p in c.text_frame.paragraphs: p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            for p in c.text_frame.paragraphs: p.font.size = Pt(11); p.font.color.rgb = DARK
    if note:
        tx = s.shapes.add_textbox(Inches(0.5), Inches(1.2 + row_h * r_count + 0.1), Inches(12), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = note; p.font.size = Pt(10); p.font.color.rgb = GRAY; p.font.italic = True

def two_col_slide(title, left_title, left_lines, right_title, right_lines):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(s, title)
    for col, (ct, lines) in enumerate([(left_title, left_lines), (right_title, right_lines)]):
        x = 0.5 + col * 6.3
        tx = s.shapes.add_textbox(Inches(x), Inches(1.1), Inches(6), Inches(0.35))
        p = tx.text_frame.paragraphs[0]; p.text = ct; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BLUE
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(6), Inches(5.3))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(30,30,46); r.line.fill.background()
        tf = r.text_frame; tf.word_wrap = True; tf.margin_top = Pt(10); tf.margin_left = Pt(12)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.font.size = Pt(11); p.font.name = 'Consolas'; p.font.color.rgb = RGBColor(205,214,244)
            p.space_before = Pt(1); p.space_after = Pt(1)


# ══════════════════════════════════════════════════════════════
# 표지
# ══════════════════════════════════════════════════════════════
title_slide("CIMS 시스템 기술 문서", "VoIP / MCPTT PTT 서버 시스템\n아키텍처 · 인터페이스 규격 · 메시지 Flow · API 명세\n\nv1.0 — 2026.03")

# ══════════════════════════════════════════════════════════════
# Part 1: 시스템 아키텍처
# ══════════════════════════════════════════════════════════════
section_slide("Part 1. 시스템 아키텍처", "컴포넌트 구성, 역할, 포트, 데이터 흐름")

# 구성도 (도형)
s = prs.slides.add_slide(prs.slide_layouts[6])
_title_bar(s, "시스템 구성도")
# 클라이언트
_box(s, 0.5, 1.3, 2.5, 1.0, "Phone App", "브라우저 (React)\nWS + WebRTC", LIGHT_BLUE)
_box(s, 3.5, 1.3, 2.5, 1.0, "Console", "브라우저 (React)\nHTTPS REST", LIGHT_BLUE)
_box(s, 6.5, 1.3, 2.5, 1.0, "VoLTE UE", "SIP 단말\nSIP + RTP", LIGHT_BLUE)
_box(s, 9.5, 1.3, 2.5, 1.0, "MCPTT UE", "PTT 단말\nSIP + RTCP Floor", LIGHT_BLUE)
# 서버
_box(s, 0.5, 3.2, 2.5, 1.5, "cwrtc", "WebSocket ↔ SIP 브리지\nDTLS-SRTP ↔ RTP\nHTTP/WS :8080", LIGHT_GREEN)
_box(s, 3.5, 3.2, 2.5, 1.5, "CSC", "REST API (Auth/Admin)\nIdMS/GMS/CMS/KMS\nHTTPS :4420/:4430", LIGHT_GREEN)
_box(s, 6.5, 3.2, 2.5, 1.5, "CSP", "SIP 시그널링 서버\n등록/통화/그룹콜\nSIP :5060", LIGHT_GREEN)
_box(s, 9.5, 3.2, 2.5, 1.5, "CMP", "RTP 미디어 릴레이\nMCPTT Floor 제어\nJSON :9000 / RTP :50000+", LIGHT_GREEN)
_box(s, 5.0, 5.6, 3.3, 0.8, "MySQL DB", "users, subscriptions\nptt_groups, call_logs", RGBColor(254,243,199))
# 연결 라벨
_arrow_text(s, 0.5, 2.4, 2.5, "▼ WebSocket + DTLS-SRTP")
_arrow_text(s, 3.5, 2.4, 2.5, "▼ HTTPS REST API")
_arrow_text(s, 6.5, 2.4, 2.5, "▼ SIP + RTP")
_arrow_text(s, 9.5, 2.4, 2.5, "▼ SIP + RTCP Floor")
_arrow_text(s, 1.5, 4.85, 4.0, "◄── SIP (UDP 5060) ──►")
_arrow_text(s, 7.5, 4.85, 3.5, "◄── UDP JSON (9000) ──►")

# 컴포넌트 상세
content_slide("CSP — Call Service Platform", [
    ("역할", ["SIP 프로토콜 서버: REGISTER, INVITE, BYE, SUBSCRIBE, NOTIFY", "1:1 VoIP 통화 및 PTT 그룹 통화 제어", "CMP에 RTP 세션 생성/삭제 명령 (UDP JSON)"]),
    ("핵심 클래스", ["CSipServer — SIP 메시지 핸들링 (Digest 인증 포함)", "CGroupCallService — PTT 그룹 통화 관리 (멤버 초대, 공유 RTP 세션)", "CmpClient — CMP에 UDP JSON 명령 전송 (타임아웃 2초)", "CSubscriptionManager — SIP SUBSCRIBE/NOTIFY 상태 관리"]),
    ("포트", ["UDP 5060 (SIP), TCP 25061, TLS 5061", "UDP 9001 (CMP 응답 수신)"]),
])

content_slide("CMP — Component Media Provider", [
    ("역할", ["RTP/RTCP 패킷 릴레이 (1:1 및 그룹)", "MCPTT Floor 제어 (RTCP APP 패킷, PT=204, name='MCPT')", "그룹 RTP 멀티캐스트: 화자→전체 멤버 포워딩"]),
    ("핵심 클래스", ["CmpServer — UDP JSON 명령 수신/디스패치, 리소스 풀 관리", "PRtpTrans — RTP/RTCP/Video 소켓 관리, 패킷 수신/전송", "McpttGroup — 그룹 멤버 관리, Floor 상태머신, SSRC/Seq 재작성"]),
    ("리소스 풀", ["세션당 4포트: Audio RTP(+0), RTCP(+1), Video RTP(+2), RTCP(+3)", "기본 20세션 × 4포트 = 50000~50079", "proc() 처리 순서: RTCP(Floor) → RTP(Audio) → Video"]),
])

content_slide("CSC — CIMS Service Controller", [
    ("역할", ["웹 인증 API (JWT, SHA256 비밀번호 해시)", "가입자/구독/그룹 CRUD (MySQL)", "MCPTT 서비스: IdMS(OAuth2 PKCE), GMS(그룹 XML), CMS(프로파일), KMS(키관리)"]),
    ("API 서버 구성", ["Admin 서버 (HTTPS :4420) — /api/v1/auth/*, /api/v1/users/*, /api/v1/ptt/groups/*", "MCPTT 서버 (HTTPS :4430) — /idms/*, /org.openmobilealliance.groups/*, /org.3gpp.mcptt.*"]),
    ("기술 스택", ["Python 3 + Uvicorn (ASGI)", "pymysql (MySQL), PyJWT (토큰), python-pptx (문서)"]),
])

content_slide("cwrtc — WebRTC-SIP Bridge", [
    ("역할", ["브라우저 WebSocket ↔ SIP 프로토콜 변환 (B2BUA)", "DTLS-SRTP(암호화) ↔ Plain RTP(CMP향) 미디어 브리지", "RTCP Floor 제어 패킷 릴레이 (브라우저 ↔ CMP)"]),
    ("포트 레이아웃 (세션당 6포트)", ["+0: Audio DTLS (브라우저↔cwrtc, 암호화)", "+1: Audio RTP (cwrtc↔CMP, 평문)", "+2: Audio RTCP (Floor 제어)", "+3: Video DTLS, +4: Video RTP, +5: Video RTCP"]),
    ("SDP 생성", ["브라우저향: UDP/TLS/RTP/SAVPF + ICE + DTLS fingerprint", "CSP향: RTP/AVP (평문, 코덱: opus/48000 또는 AMR-WB/16000)"]),
])

table_slide("포트 할당 총괄",
    ["컴포넌트", "포트", "프로토콜", "방향", "용도"],
    [
        ["CSP", "5060", "SIP/UDP", "UE↔CSP", "SIP 시그널링 (REGISTER, INVITE, BYE)"],
        ["CSP", "5061", "SIP/TLS", "UE↔CSP", "SIP 시그널링 (암호화)"],
        ["CSP", "9001", "UDP", "CSP→CMP", "CMP 제어 명령 송신"],
        ["CMP", "9000", "UDP JSON", "CSP→CMP", "RTP 세션/그룹 관리 명령"],
        ["CMP", "50000-50079", "RTP/RTCP", "UE↔CMP", "미디어 릴레이 + Floor RTCP"],
        ["CSC", "4420", "HTTPS", "Console→CSC", "인증/가입자/그룹 관리 API"],
        ["CSC", "4430", "HTTPS", "Phone→CSC", "MCPTT IdMS/GMS/CMS/KMS"],
        ["cwrtc", "8080", "HTTP/WS", "Phone→cwrtc", "WebSocket + 정적 파일 서빙"],
        ["cwrtc", "30000+", "DTLS-SRTP", "Phone↔cwrtc", "WebRTC 미디어 (암호화)"],
        ["MySQL", "3306", "TCP", "CSP/CSC→DB", "사용자/그룹/통화 데이터"],
    ])

# ══════════════════════════════════════════════════════════════
# Part 2: 메시지 Flow (상세)
# ══════════════════════════════════════════════════════════════
section_slide("Part 2. 메시지 Flow", "각 단계별 실제 메시지 및 필수 파라미터")

# --- PTT 그룹 통화 ---
code_slide("PTT 그룹 통화 Flow — ① 발신자 INVITE", """
[1] UE-A → CSP: SIP INVITE (그룹 번호로 발신)
    INVITE sip:+82571910001@csp SIP/2.0
    From: <sip:+82571900001@csp>;tag=a1b2c3
    To: <sip:+82571910001@csp>
    Call-ID: abc123@192.168.0.100
    CSeq: 1 INVITE
    Contact: <sip:+82571900001@192.168.0.100:5060>
    Content-Type: application/sdp

    v=0
    o=- 0 0 IN IP4 192.168.0.100
    s=-
    m=audio 40000 RTP/AVP 99
    c=IN IP4 192.168.0.100
    a=rtpmap:99 AMR-WB/16000
    a=sendrecv

[2] CSP → CMP: addGroup (UDP JSON)
    {"trans_id":1, "payload":{"cmd":"addgroup", "group_id":"+82571910001",
     "count":3, "members":"+82571900001:0,+821030432632:0,+82571900002:1",
     "csp_id":"CSP_MAIN", "cmp_id":"CMP_MAIN"}}

[3] CMP → CSP: 응답
    {"trans_id":1, "response":{"status":"OK","ip":"192.168.0.2","port":50000,"video_port":50002}}

[4] CSP → UE-A: 200 OK (공유 RTP 포트)
    SIP/2.0 200 OK
    m=audio 50000 RTP/AVP 99        ← CMP 공유 포트
    c=IN IP4 192.168.0.2            ← CMP IP
""", "※ group_id로 Group JSON/DB에서 멤버 조회 → members CSV로 CMP에 우선순위 전달")

code_slide("PTT 그룹 통화 Flow — ② 멤버 INVITE", """
[5] CSP → UE-B: INVITE (Multipart: SDP + MCPTT XML)
    INVITE sip:+821030432632@csp SIP/2.0
    Content-Type: multipart/mixed; boundary=mcptt
    Contact: <sip:+82571910001@csp>;isfocus      ← 그룹 포커스
    P-Called-Party-ID: <sip:+821030432632@csp>
    Answer-Mode: Auto                             ← 자동 응답 요구
    Resource-Priority: mcpttp.6

    --mcptt
    Content-Type: application/vnd.3gpp.mcptt-info+xml
    <?xml version="1.0" encoding="UTF-8"?>
    <mcpttinfo xmlns="urn:3gpp:ns:mcpttInfo:1.0">
      <mcptt-Params>
        <session-type>prearranged</session-type>
        <mcptt-request-uri>tel:+821030432632</mcptt-request-uri>
        <mcptt-calling-user-id>tel:+82571900001</mcptt-calling-user-id>
        <mcptt-calling-group-id>tel:+82571910001</mcptt-calling-group-id>
      </mcptt-Params>
    </mcpttinfo>

    --mcptt
    Content-Type: application/sdp
    m=audio 50000 RTP/AVP 99        ← 동일 공유 포트
    m=application 50001 UDP MCPTT   ← Floor 제어 포트 (RTP+1)
    a=floorid:0 mstrm:audio
    --mcptt--

[6] UE-B → CSP: 200 OK (수락, 자신의 RTP 포트 포함)
[7] CSP → CMP: joinGroup
    {"trans_id":5, "payload":{"cmd":"joingroup", "group_id":"+82571910001",
     "session_id":"+821030432632", "user_ip":"192.168.0.3",
     "user_port":10001, "user_video_port":10004}}
""")

code_slide("PTT 그룹 통화 Flow — ③ Floor 제어", """
[8] UE-A → CMP: RTCP FLOOR_REQUEST (발언권 요청)
    RTCP APP 패킷 (UDP, RTP포트+1=50001):
    ┌──────────────────────────────────────────┐
    │ V=2 | PT=204(APP) | length=5            │
    │ SSRC = 1000 (CMP 할당)                   │
    │ name = "MCPT"                            │
    │ opcode = 1 (FLOOR_REQUEST)               │
    │ id_len = 15                              │
    │ speaker_id = "+82571900001"              │
    └──────────────────────────────────────────┘

[9] CMP → UE-A: RTCP FLOOR_GRANT (발언권 승인)
    opcode = 2, ssrc = 1000, speaker_id = "+82571900001"

[10] CMP → UE-B,C: RTCP FLOOR_TAKEN (화자 알림, 브로드캐스트)
    opcode = 6, ssrc = 1000, speaker_id = "+82571900001"

[11] UE-A → CMP: RTP Audio (화자 음성)
    CMP는 화자의 RTP를 수신하여 모든 멤버에게 포워딩
    (SSRC/Sequence 재작성: 수신자별 고유값)

[12] UE-A → CMP: RTCP FLOOR_RELEASE (발언권 반납)
    opcode = 4

[13] CMP → ALL: RTCP FLOOR_IDLE (발언권 해제)
    opcode = 5
""", "※ 우선순위 선점: 현재 화자(prio=1)보다 높은 요청(prio=0) → REVOKE(7) → GRANT(2)")

# --- 웹 단말 ---
code_slide("웹 단말 연결 Flow — ① IdMS 인증 (OAuth2 PKCE)", """
[1] 브라우저 → CSC: Authorization Request
    GET https://server:4430/idms/authreq
      ?user_name=tel%3A%2B82571900001
      &user_password=123456
      &client_id=MCPTT_UE
      &redirect_uri=https%3A%2F%2Fserver%3A3000%2Fcallback
      &code_challenge=rPzO_zKmsLNK_C4NO8sZda2repcQTC-ji6mslmgR9Hc
      &code_challenge_method=S256
      &scope=openid+mcptt
      &state=omoD62dmYRyUhjp3

    응답: {"code":"3ce2dae6-5da9-4126-b587-6135d3937193", "state":"omoD62dmYRyUhjp3"}

[2] 브라우저 → CSC: Token Request
    POST https://server:4430/idms/tokenreq
    Content-Type: application/x-www-form-urlencoded

    grant_type=authorization_code
    &code=3ce2dae6-5da9-4126-b587-6135d3937193
    &redirect_uri=https://server:3000/callback
    &client_id=MCPTT_UE
    &code_verifier=test123456789012345678901234567890123456789012

    응답: {
      "access_token": "eyJhbGciOiJIUzI1NiIs...",
      "refresh_token": "a1b2c3d4-...",
      "id_token": "eyJ...",
      "expires_in": 3600
    }
""", "※ PKCE: code_challenge = Base64URL(SHA256(code_verifier)), TTL: auth_code=60s, access_token=1h")

code_slide("웹 단말 연결 Flow — ② GMS 그룹 조회", """
[3] 브라우저 → CSC: 내 그룹 목록
    GET https://server:4430/org.openmobilealliance.groups/users/tel%3A%2B82571900001
    Authorization: Bearer eyJhbGciOiJIUzI1NiIs...

    응답: [{"uri":"tel:+82571910001","display_name":"Alpha그룹","etag":"etag_1","member_count":3}]

[4] 브라우저 → CSC: 그룹 상세 (XML)
    GET .../users/tel%3A%2B82571900001/tel%3A%2B82571910001
    Accept: application/vnd.oma.poc.groups+xml

    응답:
    <?xml version="1.0" encoding="UTF-8"?>
    <resource-lists xmlns="urn:ietf:params:xml:ns:resource-lists"
      xmlns:rl="urn:ietf:params:xml:ns:resource-lists"
      xmlns:mcpttgi="urn:3gpp:ns:mcpttGroupInfo:1.0">
      <list-service uri="tel:+82571910001">
        <display-name xml:lang="en-us">Alpha그룹</display-name>
        <list>
          <entry uri="tel:+82571900001">
            <rl:display-name>테스트001</rl:display-name>
            <mcpttgi:user-priority>0</mcpttgi:user-priority>
          </entry>
          <entry uri="tel:+821030432632">
            <rl:display-name>관리자</rl:display-name>
            <mcpttgi:user-priority>0</mcpttgi:user-priority>
          </entry>
        </list>
        <mcpttgi:mcptt-video>true</mcpttgi:mcptt-video>
      </list-service>
    </resource-lists>
""")

two_col_slide("웹 단말 연결 Flow — ③ WebSocket SIP 등록 + PTT 그룹 참가",
    "SIP 등록 (register → registered)", [
        '[5] 브라우저 → cwrtc (WebSocket):',
        '{"type":"register",',
        ' "user":"+82571900001",',
        ' "password":"123456",',
        ' "domain":"csp",',
        ' "auth_id":"+82571900001"}',
        '',
        '[6] cwrtc → CSP (SIP):',
        'REGISTER sip:csp SIP/2.0',
        'From: <sip:+82571900001@csp>',
        'Authorization: Digest ...',
        '',
        '[7] CSP → cwrtc: 200 OK',
        '',
        '[8] cwrtc → 브라우저:',
        '{"type":"registered",',
        ' "user":"+82571900001"}',
    ],
    "PTT 그룹 초대 (ptt_auto_answer)", [
        '[9] CSP → cwrtc (SIP INVITE):',
        'INVITE sip:+82571900001@cwrtc',
        'Content-Type: multipart/mixed',
        '(SDP + MCPTT XML)',
        '',
        '[10] cwrtc → CSP: 200 OK (자동응답)',
        '',
        '[11] cwrtc → 브라우저 (WebSocket):',
        '{"type":"ptt_auto_answer",',
        ' "call_id":"WCSSyfqaawiJ...",',
        ' "from":"sip:+82571910001@csp",',
        ' "group_id":"+82571910001",',
        ' "sdp":"v=0\\r\\no=cwrtc ..."}',
        '',
        '[12] 브라우저 → cwrtc:',
        '{"type":"answer",',
        ' "call_id":"WCSSyfqaawiJ...",',
        ' "sdp":"v=0\\r\\no=- ..."}',
    ])

two_col_slide("웹 단말 Floor 제어 (PTT PUSH/RELEASE)",
    "PUSH (발언권 요청)", [
        '[1] 브라우저 → cwrtc (WebSocket):',
        '{"type":"ptt_request",',
        ' "call_id":"WCSSyfqaawiJ..."}',
        '',
        '[2] cwrtc → CMP (RTCP, port+1):',
        'RTCP APP: opcode=1 (REQUEST)',
        'SSRC=1000, name="MCPT"',
        '',
        '[3] CMP → cwrtc (RTCP):',
        'opcode=2 (GRANT), speaker=A',
        '',
        '[4] cwrtc → 브라우저:',
        '{"type":"ptt_floor",',
        ' "speaker":"tel:+82571900001"}',
        '',
        '[5] CMP → 다른멤버들 (RTCP):',
        'opcode=6 (TAKEN)',
        'speaker_id="+82571900001"',
        '',
        '[6] 다른멤버 cwrtc → 브라우저:',
        '{"type":"ptt_floor",',
        ' "speaker":"tel:+82571900001"}',
    ],
    "RELEASE (발언권 반납)", [
        '[1] 브라우저 → cwrtc:',
        '{"type":"ptt_release",',
        ' "call_id":"WCSSyfqaawiJ..."}',
        '',
        '[2] cwrtc → CMP (RTCP):',
        'RTCP APP: opcode=4 (RELEASE)',
        '',
        '[3] CMP → ALL (RTCP):',
        'opcode=5 (IDLE)',
        '',
        '[4] 각 cwrtc → 브라우저:',
        '{"type":"ptt_idle"}',
        '',
        '──────────────────────────',
        '발언권 거부 시:',
        '',
        '[3] CMP → 요청자 (RTCP):',
        'opcode=3 (REJECT)',
        '',
        '[4] cwrtc → 브라우저:',
        '{"type":"ptt_reject"}',
        '',
        '※ 우선순위 낮은 경우 거부',
    ])

code_slide("1:1 VoIP 통화 Flow (상세)", """
[1] 브라우저-A → cwrtc-A (WebSocket):
    {"type":"call", "to":"+821030432632",
     "sdp":"v=0\\r\\no=- 0 0 IN IP4 0.0.0.0\\r\\n..."}  ← WebRTC offer SDP

[2] cwrtc-A → CSP: SIP INVITE (Plain RTP SDP로 변환)
    INVITE sip:+821030432632@csp SIP/2.0
    m=audio 30001 RTP/AVP 111        ← cwrtc RTP 포트
    a=rtpmap:111 opus/48000/2

[3] CSP → CMP: add (발신자 RTP 세션)
    {"payload":{"cmd":"add","session_id":"call_001","remote_ip":"192.168.0.2","remote_port":30001}}
    응답: {"status":"OK","local_ip":"192.168.0.2","local_port":50000}

[4] CSP → cwrtc-B: SIP INVITE (CMP 릴레이 포트)
    m=audio 50000 RTP/AVP 111
    c=IN IP4 192.168.0.2

[5] cwrtc-B → 브라우저-B: {"type":"incoming","call_id":"...","from":"+82571900001","sdp":"..."}
[6] 브라우저-B → cwrtc-B: {"type":"answer","call_id":"...","sdp":"..."}  ← WebRTC answer
[7] cwrtc-B → CSP: 200 OK (RTP 포트 포함)
[8] CSP → CMP: add (착신자 RTP 세션)

[9] 미디어 흐름:
    브라우저-A ══DTLS══► cwrtc-A ══RTP══► CMP ══RTP══► cwrtc-B ══DTLS══► 브라우저-B
""")

# ══════════════════════════════════════════════════════════════
# Part 3: 인터페이스 규격 (상세)
# ══════════════════════════════════════════════════════════════
section_slide("Part 3. 인터페이스 규격", "CSP↔CMP JSON, WebSocket, RTCP Floor, SIP 메시지 샘플")

two_col_slide("CSP→CMP JSON 프로토콜 — 요청/응답 샘플",
    "addgroup 요청", [
        '// 전송: CSP → CMP (UDP, port 9000)',
        '{',
        '  "trans_id": 1,',
        '  "payload": {',
        '    "cmd": "addgroup",',
        '    "group_id": "+82571910001",',
        '    "count": 3,',
        '    "members": "+82571900001:0,',
        '               +821030432632:0,',
        '               +82571900002:1",',
        '    "csp_id": "CSP_MAIN",',
        '    "csp_sess_id": "0",',
        '    "cmp_id": "CMP_MAIN",',
        '    "cmp_sess_id": "0"',
        '  }',
        '}',
        '',
        '// members 형식: "userId:priority,..."',
        '// priority: 0=최고, 값이 클수록 낮음',
    ],
    "응답 + joingroup 요청", [
        '// addgroup 응답:',
        '{',
        '  "trans_id": 1,',
        '  "response": {',
        '    "status": "OK",',
        '    "ip": "192.168.0.2",',
        '    "port": 50000,',
        '    "video_port": 50002',
        '  }',
        '}',
        '',
        '// joingroup 요청:',
        '{',
        '  "trans_id": 5,',
        '  "payload": {',
        '    "cmd": "joingroup",',
        '    "group_id": "+82571910001",',
        '    "session_id": "+821030432632",',
        '    "user_ip": "192.168.0.3",',
        '    "user_port": 10001,',
        '    "user_video_port": 10004',
        '  }',
        '}',
    ])

table_slide("CSP→CMP 전체 명령 명세",
    ["명령", "필수 파라미터", "응답 필드", "설명"],
    [
        ["add", "session_id, remote_ip, remote_port,\nremote_video_port, peer_index", "local_ip, local_port,\nlocal_video_port", "1:1 RTP 릴레이 세션 생성/수정"],
        ["remove", "session_id", "OK", "RTP 세션 해제, 포트 반환"],
        ["addgroup", "group_id, count,\nmembers (CSV: id:prio,...)", "ip, port, video_port", "PTT 그룹 생성 + 공유 세션 할당"],
        ["joingroup", "group_id, session_id,\nuser_ip, user_port, user_video_port", "OK", "멤버를 그룹에 추가 (IP:Port 등록)"],
        ["leavegroup", "group_id, session_id", "OK", "멤버를 그룹에서 제거"],
        ["removegroup", "group_id", "OK", "그룹 + 공유 세션 삭제"],
        ["modifygroup", "group_id, members (CSV)", "OK", "멤버 우선순위 갱신"],
        ["alive", "(없음)", "OK", "연결 확인 (3초 간격 KeepAlive)"],
    ], "※ 공통 헤더: csp_id, csp_sess_id, cmp_id, cmp_sess_id / 타임아웃: 2000ms")

table_slide("RTCP APP Floor Control 패킷 상세",
    ["Opcode", "이름", "방향", "SSRC", "speaker_id", "설명"],
    [
        ["1", "FLOOR_REQUEST", "UE→CMP", "CMP 할당값", "요청자 ID", "발언권 요청 (PTT 누름)"],
        ["2", "FLOOR_GRANT", "CMP→요청자", "요청자 SSRC", "요청자 ID", "발언권 승인 → RTP 송신 시작"],
        ["3", "FLOOR_REJECT", "CMP→요청자", "요청자 SSRC", "요청자 ID", "발언권 거부 (현재 화자 우선순위 높음)"],
        ["4", "FLOOR_RELEASE", "UE→CMP", "CMP 할당값", "(없음)", "발언권 반납 (PTT 해제)"],
        ["5", "FLOOR_IDLE", "CMP→전체", "0", "(없음)", "발언권 해제됨 (아무도 안 잡고 있음)"],
        ["6", "FLOOR_TAKEN", "CMP→전체", "화자 SSRC", "화자 ID", "발언권 점유됨 (화자 정보 포함)"],
        ["7", "FLOOR_REVOKE", "CMP→현재화자", "현재화자 SSRC", "현재화자 ID", "강제 회수 (높은 우선순위 선점)"],
    ], "※ 패킷: PT=204(APP), name='MCPT', RTCP 포트=RTP포트+1, speaker_id는 4바이트 정렬")

two_col_slide("cwrtc WebSocket 프로토콜 — 클라이언트→서버",
    "register (SIP 등록)", [
        '{"type": "register",',
        ' "user": "+82571900001",',
        ' "password": "123456",',
        ' "domain": "csp",',
        ' "auth_id": "+82571900001"}',
        '',
        '// auth_id: SIP IMPI (미입력시 user와 동일)',
        '// domain: CSP의 realm 설정값',
    ],
    "call / answer / hangup", [
        '// 발신:',
        '{"type": "call",',
        ' "to": "+821030432632",',
        ' "sdp": "v=0\\r\\n..."}',
        '',
        '// 착신 수락:',
        '{"type": "answer",',
        ' "call_id": "abc123...",',
        ' "sdp": "v=0\\r\\n..."}',
        '',
        '// 통화 종료:',
        '{"type": "hangup",',
        ' "call_id": "abc123..."}',
        '',
        '// PTT 발언권:',
        '{"type": "ptt_request",',
        ' "call_id": "abc123..."}',
        '{"type": "ptt_release",',
        ' "call_id": "abc123..."}',
    ])

two_col_slide("cwrtc WebSocket 프로토콜 — 서버→클라이언트",
    "등록/통화 이벤트", [
        '// 등록 성공:',
        '{"type": "registered",',
        ' "user": "+82571900001"}',
        '',
        '// 일반 착신:',
        '{"type": "incoming",',
        ' "call_id": "abc...",',
        ' "from": "+821030432632",',
        ' "sdp": "v=0...",',
        ' "ptt": "false"}',
        '',
        '// PTT 그룹 초대 (자동 응답):',
        '{"type": "ptt_auto_answer",',
        ' "call_id": "WCSab...",',
        ' "from": "sip:+82571910001@csp",',
        ' "group_id": "+82571910001",',
        ' "sdp": "v=0..."}',
        '',
        '// 통화 종료:',
        '{"type": "ended",',
        ' "reason": "normal"}',
    ],
    "PTT Floor 이벤트", [
        '// 화자 변경 (발언권 승인):',
        '{"type": "ptt_floor",',
        ' "speaker": "tel:+82571900001"}',
        '',
        '// 발언권 해제:',
        '{"type": "ptt_idle"}',
        '',
        '// 발언권 거부:',
        '{"type": "ptt_reject"}',
        '',
        '// 멤버 참가:',
        '{"type": "ptt_member_joined",',
        ' "user_id": "tel:+82571900002"}',
        '',
        '// 멤버 퇴장:',
        '{"type": "ptt_member_left",',
        ' "user_id": "tel:+82571900002"}',
        '',
        '// 연결 유지 (30초 간격):',
        '{"type": "ping"}',
    ])

# ══════════════════════════════════════════════════════════════
# Part 4: 관리 API 명세 (상세)
# ══════════════════════════════════════════════════════════════
section_slide("Part 4. 관리 API 명세", "인증, 사용자, 구독, 그룹 — 실제 요청/응답 JSON 포함")

two_col_slide("인증 API — 로그인 / 회원가입",
    "POST /api/v1/auth/login", [
        '// 요청:',
        'POST https://server:4420/api/v1/auth/login',
        'Content-Type: application/json',
        '',
        '{"login_id": "admin",',
        ' "password": "1234"}',
        '',
        '// 성공 응답 (200):',
        '{"token": "eyJhbGciOiJIUzI1...",',
        ' "user": {',
        '   "id": 33,',
        '   "name": "관리자",',
        '   "login_id": "admin",',
        '   "role": "admin",',
        '   "ptt_subscriptions": [{',
        '     "id": "+821030432632",',
        '     "auth_id": "user@domain",',
        '     "passwd": "1234",',
        '     "dnd": false,',
        '     "forward_id": ""',
        '   }]}}',
    ],
    "POST /api/v1/auth/register", [
        '// 요청:',
        '{"name": "홍길동",',
        ' "login_id": "hong001",',
        ' "password": "mypassword"}',
        '',
        '// 성공 응답 (201):',
        '{"token": "eyJ...",',
        ' "user": {',
        '   "id": 34,',
        '   "name": "홍길동",',
        '   "login_id": "hong001",',
        '   "role": "user",',
        '   "call_subscriptions": [],',
        '   "ptt_subscriptions": []}}',
        '',
        '// 에러 응답:',
        '400: {"error":"아이디를 입력하세요"}',
        '409: {"error":"이미 사용 중인 아이디"}',
        '',
        '// JWT Payload:',
        '{"sub":"33", "login_id":"admin",',
        ' "role":"admin", "exp":17755...}',
        '// TTL: 7일, Algorithm: HS256',
    ])

two_col_slide("사용자 CRUD — 목록 조회 / 생성",
    "GET /api/v1/users (전체 목록)", [
        '// 요청:',
        'GET /api/v1/users',
        'Authorization: Bearer <admin_token>',
        '',
        '// 응답 (200):',
        '[{',
        '  "id": 1,',
        '  "name": "테스트001",',
        '  "login_id": "test001",',
        '  "org_id": "본부1팀",',
        '  "details": null,',
        '  "reject_id": [],',
        '  "call_subscriptions": [{',
        '    "id": "+821012345678",',
        '    "auth_id": "+821012345678",',
        '    "dnd": false,',
        '    "forward_id": "",',
        '    "register_time": "2026-03...",',
        '    "logout_time": null}],',
        '  "ptt_subscriptions": [...]',
        '}]',
    ],
    "POST /api/v1/users (사용자 생성)", [
        '// 요청:',
        'POST /api/v1/users',
        'Authorization: Bearer <admin_token>',
        'Content-Type: application/json',
        '',
        '{"name": "신규사용자",',
        ' "login_id": "user001",',
        ' "org_id": "본부1팀",',
        ' "details": "비고 내용",',
        ' "reject_id": ["+821099990001"]}',
        '',
        '// 성공 응답 (201):',
        '{"id": 35, "name": "신규사용자",',
        ' "login_id": "user001", ...}',
        '',
        '// 수정: PUT /api/v1/users/35',
        '{"name": "변경이름"}',
        '',
        '// 삭제: DELETE /api/v1/users/35',
        '→ 연결된 Call/PTT 구독도 삭제',
    ])

two_col_slide("Call/PTT 구독 관리",
    "POST /users/{pid}/call (Call 번호 추가)", [
        '// 요청:',
        'POST /api/v1/users/1/call',
        '',
        '{',
        '  "id": "+821012345678",',
        '  "auth_id": "+821012345678",',
        '  "passwd": "1234",',
        '  "dnd": false,',
        '  "forward_id": ""',
        '}',
        '',
        '// 필드 설명:',
        '// id: MSISDN (E.164, +국가코드)',
        '// auth_id: SIP Digest 인증 ID',
        '//          (미입력시 id와 동일)',
        '// passwd: SIP Digest 비밀번호',
        '// dnd: 방해금지 모드',
        '// forward_id: 착신전환 번호',
    ],
    "POST /users/{pid}/ptt (PTT 번호 추가)", [
        '// 요청:',
        'POST /api/v1/users/1/ptt',
        '',
        '{',
        '  "id": "+82571900001",',
        '  "auth_id":',
        '    "4503382571900001@ptt.mnc033',
        '     .mcc450.3gppnetwork.org",',
        '  "passwd": "123456",',
        '  "dnd": false,',
        '  "forward_id": ""',
        '}',
        '',
        '// auth_id 형식 (3GPP IMPI):',
        '// <MCC><MNC><MSISDN>@ptt.mnc',
        '//   <MNC>.mcc<MCC>.3gppnetwork.org',
        '',
        '// 수정: PUT .../ptt/+82571900001',
        '// 삭제: DELETE .../ptt/+82571900001',
    ])

two_col_slide("PTT 그룹 관리",
    "POST /api/v1/ptt/groups (그룹 생성)", [
        '// 요청:',
        '{',
        '  "id": "+82571910001",',
        '  "name": "Alpha그룹",',
        '  "video_enabled": true,',
        '  "members": [',
        '    {"user_id":"+82571900001",',
        '     "priority": 0},',
        '    {"user_id":"+821030432632",',
        '     "priority": 0},',
        '    {"user_id":"+82571900002",',
        '     "priority": 1}',
        '  ]',
        '}',
        '',
        '// priority: 0=최고 우선순위',
        '// video_enabled: 영상 PTT 지원',
    ],
    "그룹 조회 / 멤버 관리", [
        '// 전체 목록: GET /api/v1/ptt/groups',
        '[{"id":"+82571910001",',
        '  "name":"Alpha그룹",',
        '  "video_enabled": true,',
        '  "member_count": 3,',
        '  "members":[',
        '    {"user_id":"+82571900001",',
        '     "priority":0}, ...]}]',
        '',
        '// 멤버 추가:',
        'POST /api/v1/ptt/groups/',
        '  +82571910001/members',
        '{"user_id":"+82571900003",',
        ' "priority": 2}',
        '',
        '// 멤버 삭제:',
        'DELETE /api/v1/ptt/groups/',
        '  +82571910001/members/',
        '  +82571900003',
    ])

table_slide("DB 스키마 상세",
    ["테이블", "컬럼", "타입", "키", "설명"],
    [
        ["users", "id", "INT AUTO", "PK", "사용자 고유 ID"],
        ["", "name", "VARCHAR(128)", "", "표시 이름"],
        ["", "login_id", "VARCHAR(255)", "UNIQUE", "로그인 아이디 (ex: test001)"],
        ["", "password", "VARCHAR(255)", "", "SHA256 해시 비밀번호"],
        ["", "role", "ENUM(admin,user)", "", "역할"],
        ["", "org_id", "VARCHAR(64)", "", "소속 조직"],
        ["voip_subscriptions", "id", "VARCHAR(32)", "PK", "MSISDN (E.164)"],
        ["", "user_id", "INT", "FK→users", "소유자"],
        ["", "auth_id", "VARCHAR(255)", "", "SIP IMPI"],
        ["", "passwd", "VARCHAR(64)", "", "SIP Digest PW"],
        ["", "dnd", "TINYINT", "", "방해금지 (0/1)"],
        ["", "forward_id", "VARCHAR(32)", "", "착신전환 번호"],
        ["ptt_groups", "id", "VARCHAR(32)", "PK", "그룹 MSISDN"],
        ["", "name", "VARCHAR(128)", "", "그룹 표시 이름"],
        ["", "video_enabled", "TINYINT", "", "영상 지원 여부"],
        ["ptt_group_members", "group_id", "VARCHAR(32)", "FK", "그룹 ID"],
        ["", "user_id", "VARCHAR(32)", "FK", "PTT 구독 MSISDN"],
        ["", "priority", "INT", "", "우선순위 (0=최고)"],
    ])

# 마지막
title_slide("감사합니다", "CIMS Technical Documentation v1.0\n\n3GPP TS 24.379/380/481 · OMA PoC · RFC 3261/3550")

output = '/home/nex/work/cims/docs/CIMS_Technical_Document.pptx'
prs.save(output)
print(f'PPT 생성 완료: {output} ({len(prs.slides)}페이지)')
